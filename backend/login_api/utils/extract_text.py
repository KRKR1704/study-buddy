# utils/extract_text.py
import os
from typing import List

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

try:
    from striprtf.striprtf import rtf_to_text
except Exception:
    rtf_to_text = None

try:
    from bs4 import BeautifulSoup
except Exception:
    BeautifulSoup = None  # type: ignore

try:
    import pandas as pd
except Exception:
    pd = None  # type: ignore

try:
    from ebooklib import epub
except Exception:
    epub = None  # type: ignore


def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[-1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext == ".pptx":
        return _extract_pptx(file_path)
    if ext in (".txt", ".md", ".markdown"):
        return _extract_plain(file_path)
    if ext == ".rtf":
        return _extract_rtf(file_path)
    if ext in (".html", ".htm"):
        return _extract_html(file_path)
    if ext == ".csv":
        return _extract_csv(file_path)
    if ext == ".xlsx":
        return _extract_xlsx(file_path)
    if ext == ".epub":
        return _extract_epub(file_path)

    # Explicit message for legacy Office formats
    if ext in (".doc", ".ppt"):
        raise ValueError(
            "Legacy Office format (.doc/.ppt) is not supported. "
            "Please convert to DOCX/PPTX and upload again."
        )

    raise ValueError(
        "Unsupported file format. Please upload a text document "
        "(PDF/DOCX/PPTX/TXT/RTF/MD/HTML/CSV/XLSX/EPUB)."
    )


def _extract_pdf(path: str) -> str:
    with fitz.open(path) as doc:
        return "\n".join([page.get_text() for page in doc])


def _extract_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def _extract_pptx(path: str) -> str:
    prs = Presentation(path)
    lines: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def _extract_plain(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_rtf(path: str) -> str:
    if not rtf_to_text:
        raise ValueError("RTF support requires 'striprtf'. Please install it.")
    with open(path, "rb") as f:
        data = f.read().decode("utf-8", errors="ignore")
    return rtf_to_text(data)


def _extract_html(path: str) -> str:
    if not BeautifulSoup:
        raise ValueError("HTML support requires 'beautifulsoup4'. Please install it.")
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "lxml") if _has_lxml() else BeautifulSoup(html, "html.parser")
    return soup.get_text(separator="\n").strip()


def _extract_csv(path: str) -> str:
    if not pd:
        raise ValueError("CSV support requires 'pandas'. Please install it.")
    df = pd.read_csv(path, dtype=str, encoding="utf-8", engine="python")
    df = df.fillna("")
    return "\n".join(df.apply(lambda row: " | ".join(map(str, row.values)), axis=1).tolist())


def _extract_xlsx(path: str) -> str:
    if not pd:
        raise ValueError("XLSX support requires 'pandas' and 'openpyxl'. Please install them.")
    xls = pd.ExcelFile(path)
    lines: List[str] = []
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name, dtype=str)
        df = df.fillna("")
        lines.append(f"=== Sheet: {sheet_name} ===")
        lines.extend(df.apply(lambda row: " | ".join(map(str, row.values)), axis=1).tolist())
    return "\n".join(lines)


def _extract_epub(path: str) -> str:
    if not epub:
        raise ValueError("EPUB support requires 'ebooklib'. Please install it.")
    book = epub.read_epub(path)
    result: List[str] = []
    for item in book.get_items():
        if item.get_type() == 9:  # DOCUMENT
            try:
                content = item.get_body_content().decode("utf-8", errors="ignore")
                if BeautifulSoup:
                    soup = BeautifulSoup(content, "lxml") if _has_lxml() else BeautifulSoup(content, "html.parser")
                    text = soup.get_text(separator="\n")
                else:
                    text = content
                if text.strip():
                    result.append(text.strip())
            except Exception:
                continue
    return "\n\n".join(result)


def _has_lxml() -> bool:
    try:
        import lxml  # noqa
        return True
    except Exception:
        return False
